#!/usr/bin/env python3
"""Update PPTX with new slides and content while preserving visual style."""
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN

INPUT  = 'AI-powered assessment platform for higher education (1).pptx'
OUTPUT = 'AI-powered assessment platform for higher education (updated).pptx'

# ── colour helpers ────────────────────────────────────────────────────────────
WHITE  = RGBColor(0xF2, 0xF2, 0xF2)
PWHITE = RGBColor(0xFF, 0xFF, 0xFF)
CYAN   = RGBColor(0x77, 0xE8, 0xFF)
PINK   = RGBColor(0xF5, 0x9A, 0xF4)
BG     = RGBColor(0x49, 0x16, 0x86)

# ── xml helpers ───────────────────────────────────────────────────────────────

def _make_rPr(size_pt, bold, color: RGBColor, font_name):
    rPr = etree.Element(qn('a:rPr'))
    rPr.set('lang', 'en-US')
    rPr.set('dirty', '0')
    if size_pt:
        rPr.set('sz', str(int(size_pt * 100)))
    if bold:
        rPr.set('b', '1')
    if color:
        sf = etree.SubElement(rPr, qn('a:solidFill'))
        sc = etree.SubElement(sf, qn('a:srgbClr'))
        sc.set('val', str(color).upper())
    if font_name:
        lat = etree.SubElement(rPr, qn('a:latin'))
        lat.set('typeface', font_name)
    return rPr


def _make_run(text, size_pt, bold, color: RGBColor, font_name):
    r = etree.Element(qn('a:r'))
    r.append(_make_rPr(size_pt, bold, color, font_name))
    t = etree.SubElement(r, qn('a:t'))
    t.text = text
    return r


def _make_para(runs, align=None):
    p = etree.Element(qn('a:p'))
    if align:
        pPr = etree.SubElement(p, qn('a:pPr'))
        pPr.set('algn', align)
    for r in runs:
        p.append(r)
    return p


def replace_txBody(shape, paragraphs):
    """Replace content of a text-frame shape.
    paragraphs: list of lists of (text, size_pt, bold, color, font_name)
    Each inner list = one paragraph; empty inner list = blank line.
    """
    txBody = shape.text_frame._txBody
    # remove existing paragraphs
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)
    for para in paragraphs:
        if not para:
            txBody.append(etree.Element(qn('a:p')))
        else:
            runs = [_make_run(t, sz, b, c, f) for t, sz, b, c, f in para]
            txBody.append(_make_para(runs))


# ── slide utilities ───────────────────────────────────────────────────────────

def clone_slide(prs, src_idx):
    """Deep-copy slide[src_idx] and append it at the end; return new slide."""
    src = prs.slides[src_idx]
    new_slide = prs.slides.add_slide(src.slide_layout)
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        sp_tree.remove(child)
    for child in src.shapes._spTree:
        sp_tree.append(copy.deepcopy(child))
    try:
        src_color = src.background.fill.fore_color.rgb
        new_slide.background.fill.solid()
        new_slide.background.fill.fore_color.rgb = src_color
    except Exception:
        pass
    return new_slide


def move_slide(prs, old_idx, new_idx):
    xml_slides = prs.slides._sldIdLst
    sl = list(xml_slides)
    el = sl[old_idx]
    xml_slides.remove(el)
    xml_slides.insert(new_idx, el)


def add_textbox(slide, left, top, width, height,
                paragraphs, word_wrap=False):
    """Add a new textbox to slide and fill with paragraphs.
    paragraphs: same format as replace_txBody.
    """
    from pptx.util import Emu
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    txBox.text_frame.word_wrap = word_wrap
    replace_txBody(txBox, paragraphs)
    return txBox


def get_text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame]


# ── table helper ──────────────────────────────────────────────────────────────

def add_styled_table(slide, left, top, width, height, headers, rows,
                     col_widths=None):
    """Add a table styled with dark-purple theme."""
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    tbl = slide.shapes.add_table(n_rows, n_cols,
                                  Emu(left), Emu(top),
                                  Emu(width), Emu(height)).table

    # Column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Emu(w)

    def style_cell(cell, text, is_header=False, is_bold=False, fg=None):
        cell.text = ''
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x35, 0x0C, 0x6B) if is_header else RGBColor(0x3D, 0x10, 0x78)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(11) if not is_header else Pt(12)
        run.font.bold = is_header or is_bold
        run.font.color.rgb = CYAN if is_header else (PWHITE if text in ('✓','✗') else WHITE)
        run.font.name = 'Cairo'

    # Header row
    for j, h in enumerate(headers):
        style_cell(tbl.cell(0, j), h, is_header=True)

    # Data rows
    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            style_cell(tbl.cell(i + 1, j), cell_text,
                       is_bold=(j == 0))

    return tbl


# ═══════════════════════════════════════════════════════════════════════════════
# Load presentation
# ═══════════════════════════════════════════════════════════════════════════════
prs = Presentation(INPUT)
W = prs.slide_width   # 18288000
H = prs.slide_height  # 10287000

# Short aliases to original slides (0-based)
S = prs.slides

# ═══════════════════════════════════════════════════════════════════════════════
# STEP 1 – Modify existing slides
# ═══════════════════════════════════════════════════════════════════════════════

# ── Slide 1 (Title) – add subtitle info ──────────────────────────────────────
slide1 = S[0]
# Add subtitle textbox at bottom-left
add_textbox(slide1,
    left=503564, top=6994351, width=5040548, height=2959735,
    paragraphs=[
        [('Graduation Project', 13, False, WHITE, 'Cairo')],
        [],
        [('Supervisor: ', 11, True, CYAN, 'Cairo'), ('Laid Kahloul', 11, False, WHITE, 'Cairo')],
        [('Co-supervisor: ', 11, True, CYAN, 'Cairo'), ('Nour El Houda Ben Chaabene', 11, False, WHITE, 'Cairo')],
        [('By: ', 11, True, CYAN, 'Cairo'), ('MERIEM CHEBANA', 11, False, WHITE, 'Cairo')],
        [],
        [('2025 / 2026', 11, False, WHITE, 'Cairo')],
    ],
    word_wrap=True)

# ── Slide 2 (Contents) – update section list ─────────────────────────────────
slide2 = S[1]
tshapes2 = get_text_shapes(slide2)
# tshapes2[0] = "Contents" title, tshapes2[1] = items list
replace_txBody(tshapes2[1], [
    [('01 · Problematic & Objectives', 13, True, WHITE, 'Cairo Semi-Bold')],
    [('       — Three Dimensions  |  Research Gap & System Scope  |  Objectives', 10, False, CYAN, 'Cairo')],
    [],
    [(' 02 · Related Work & Inspirations', 13, True, WHITE, 'Cairo Semi-Bold')],
    [('       — General Studies  |  Recent Key Studies  |  System Positioning  |  Shared Pillars', 10, False, CYAN, 'Cairo')],
    [],
    [(' 03 · Proposition & Contribution', 13, True, WHITE, 'Cairo Semi-Bold')],
    [('       — Contribution  |  System Performance & Validation Results', 10, False, CYAN, 'Cairo')],
    [],
    [(' 04 · System Architecture', 13, True, WHITE, 'Cairo Semi-Bold')],
    [('       — Four-Layer Architecture  |  docParser Pipeline  |  Indexing  |  Generation  |  Evaluation  |  Web App', 10, False, CYAN, 'Cairo')],
    [],
    [(' 05 · Problems Encountered & Improvements', 13, True, WHITE, 'Cairo Semi-Bold')],
    [],
    [(' 06 · Results & System Performance', 13, True, WHITE, 'Cairo Semi-Bold')],
    [],
    [(' 07 · Models & Technical Stack', 13, True, WHITE, 'Cairo Semi-Bold')],
])

# ── Slide 4 (Problematic) – add "Three Dimensions" label + bilingual line ────
slide4 = S[3]
tshapes4 = get_text_shapes(slide4)
# tshapes4 that has the content text (the big one with the 3 bullet points)
content_box = max([s for s in tshapes4 if s.text_frame.text.strip().startswith('Time') or
                   s.text_frame.text.strip().startswith(' ')],
                   key=lambda s: len(s.text_frame.text), default=None)
if content_box is None:
    content_box = tshapes4[-1]

replace_txBody(content_box, [
    [('Time & resources: ', 10, True, CYAN, 'Cairo'),
     ('Preparing balanced university exams consumes an enormous amount of a professor\'s time. Beyond the time cost, there is a constant risk of producing unbalanced assessments that fail to accurately measure student competencies, or of relying on repetitive templates that limit pedagogical creativity.', 10, False, WHITE, 'Cairo')],
    [],
    [('Open-ended question difficulty: ', 10, True, CYAN, 'Cairo'),
     ('Constructed-response questions are expensive to design and lack standardized answer formats, making automated or consistent evaluation difficult.', 10, False, WHITE, 'Cairo')],
    [],
    [('Academic integrity: ', 10, True, CYAN, 'Cairo'),
     ('The need to produce multiple equivalent exam versions (to reduce cheating) amplifies all of these pressures. In the Algerian university context, this challenge is compounded by bilingual course archives in French and Arabic — formats that existing automated tools cannot handle without dedicated language-aware processing.', 10, False, WHITE, 'Cairo')],
])

# Also update title box to say "Three Dimensions — Problematic"
title_box4 = [s for s in tshapes4 if 'Problematic' in s.text_frame.text][0]
replace_txBody(title_box4, [
    [('Three Dimensions — Problematic', 30, False, WHITE, 'Zen Dots')],
])

# ── Slide 5 (Objectives) – add Bilingual support ─────────────────────────────
slide5 = S[4]
tshapes5 = get_text_shapes(slide5)
obj_box = max(tshapes5, key=lambda s: len(s.text_frame.text))
replace_txBody(obj_box, [
    [(' Automate exam preparation', 13, True, PINK, 'Cairo Bold')],
    [('       Generate high-quality questions automatically from the professor\'s own academic materials.', 10, False, WHITE, 'Cairo')],
    [],
    [(' Pedagogical alignment', 13, True, PINK, 'Cairo Bold')],
    [('       Analyze lectures, slides, and past exams to ensure questions match intended learning outcomes (Bloom\'s taxonomy).', 10, False, WHITE, 'Cairo')],
    [],
    [(' Co-creation tool', 13, True, PINK, 'Cairo Bold')],
    [('        Act as an assistant — the professor reviews, edits, and approves before any exam is finalized.', 10, False, WHITE, 'Cairo')],
    [],
    [(' Transparency & metadata', 13, True, PINK, 'Cairo Bold')],
    [('        Every question carries: difficulty level, targeted Bloom\'s level, and solutions.', 10, False, WHITE, 'Cairo')],
    [],
    [(' Bilingual support', 13, True, PINK, 'Cairo Bold')],
    [('        Handle French and Arabic course materials natively, without translation, routing each language to its appropriate retrieval model.', 10, False, WHITE, 'Cairo')],
])

# ── Slide 6 (Related Work) – restructure with named studies ──────────────────
slide6 = S[5]
tshapes6 = get_text_shapes(slide6)
content6 = max(tshapes6, key=lambda s: len(s.text_frame.text))
replace_txBody(content6, [
    [('iTELL', 10, True, CYAN, 'Cairo'),
     (' — Automated question generation and evaluation of constructed responses using GPT-3.5, MPNet and BLEURT with consensus voting. Accuracy: 0.81. Emphasized the importance of open-ended questions (Generation Effect).', 10, False, WHITE, 'Cairo')],
    [],
    [('ItemForge', 10, True, CYAN, 'Cairo'),
     (' — Mathematics assessment items using Anderson & Krathwohl taxonomy with a multi-agent approach (generation + review). Accuracy: ~0.8.', 10, False, WHITE, 'Cairo')],
    [],
    [('Institutional Regulations Study', 10, True, CYAN, 'Cairo'),
     (' — MCQ generation from large PDFs evaluated with GPT-4o, reaching 89.9% accuracy. Prompt engineering and self-review significantly improved quality.', 10, False, WHITE, 'Cairo')],
    [],
    [('Large-Scale Field Study (1,686 students)', 10, True, CYAN, 'Cairo'),
     (' — AI-generated exams can match or outperform human-written ones using an iterative refinement strategy.', 10, False, WHITE, 'Cairo')],
])

# ── Slide 7 (Shared Pillars) – add Mucciaccia and Meissner references ─────────
slide7 = S[6]
tshapes7 = get_text_shapes(slide7)
content7 = max(tshapes7, key=lambda s: len(s.text_frame.text))
replace_txBody(content7, [
    [('From ItemForge', 10, True, CYAN, 'Cairo'),
     (': use of Anderson & Krathwohl taxonomy to control cognitive levels, and the multi-agent concept — generation separated from evaluation — implemented here with multiple LLMs across pipeline stages.', 10, False, WHITE, 'Cairo')],
    [],
    [('From the Institutional Regulations Study', 10, True, CYAN, 'Cairo'),
     (': prompt engineering and a self-review mechanism (AI-Judge) to improve question quality and reduce errors.', 10, False, WHITE, 'Cairo')],
    [],
    [('From the Large-Scale Field Study', 10, True, CYAN, 'Cairo'),
     (': iterative refinement to progressively improve the quality of outputs.', 10, False, WHITE, 'Cairo')],
    [],
    [('From Mucciaccia et al. (2025)', 10, True, CYAN, 'Cairo'),
     (': schema enforcement and LLM-as-judge — extended here with Pydantic + auto-retry (max 3 retries).', 10, False, WHITE, 'Cairo')],
    [],
    [('From Meissner et al. (2024)', 10, True, CYAN, 'Cairo'),
     (': Bloom classification must be done by an independent model, not the generator itself.', 10, False, WHITE, 'Cairo')],
])

# ── Slide 16 (Post-Gen Eval) – update 3 criteria → 4 criteria ────────────────
slide16 = S[15]
tshapes16 = get_text_shapes(slide16)
content16 = max(tshapes16, key=lambda s: len(s.text_frame.text))
replace_txBody(content16, [
    [('The goal is to verify the quality of each generated question before delivery.', 9, False, WHITE, 'Cairo')],
    [],
    [('Re-Retrieval', 9, False, CYAN, 'Cairo'),
     ('  The question text (truncated to 300 characters) is used as a new RAG query, eliminating circular evaluation bias.', 9, False, WHITE, 'Cairo')],
    [],
    [('Question Formatting', 9, False, CYAN, 'Cairo'),
     ('  The question is reformatted for evaluation including the correct answer (MCQ) and retrieved context (up to 2,500 chars).', 9, False, WHITE, 'Cairo')],
    [],
    [('Evaluation Axes', 9, False, CYAN, 'Cairo'),
     ('  The question is evaluated on four criteria:', 9, False, WHITE, 'Cairo')],
    [('      (1) Answer Extractability — can the answer be found in the source material?', 9, False, WHITE, 'Cairo')],
    [('      (2) Cognitive Level Alignment — does it match the required Bloom\'s level?', 9, False, WHITE, 'Cairo')],
    [('      (3) Question Quality — good / needs revision / reject', 9, False, WHITE, 'Cairo')],
    [('      (4) Difficulty Fairness — is the difficulty appropriate for the level?', 9, False, WHITE, 'Cairo')],
    [],
    [('Decision Logic', 9, False, CYAN, 'Cairo'),
     ('  Pass: facts correct and quality acceptable. Flag: incorrect level or needs revision. Reject: incorrect facts or poor quality (priority over flag).', 9, False, WHITE, 'Cairo')],
    [],
    [('Silent Fallback', 9, False, CYAN, 'Cairo'),
     ('  If evaluation fails (connection error), a default pass is returned with a note that evaluation was unavailable.', 9, False, WHITE, 'Cairo')],
])

# ── Slide 18 (Contribution) – 3-Axis → 4-Axis ────────────────────────────────
slide18 = S[17]
tshapes18 = get_text_shapes(slide18)
content18 = [s for s in tshapes18 if 'Hybrid' in s.text_frame.text][0]
replace_txBody(content18, [
    [('Hybrid Retrieval — RRF Fusion', 10, False, CYAN, 'Cairo')],
    [(' The system combines three retrievers: BM25 (keywords), TF-IDF + LSA (semantic, no GPU), and AraBERT (Arabic). Fused using Reciprocal Rank Fusion with k = 60.', 10, False, WHITE, 'Cairo')],
    [],
    [('Native Multilingualism', 10, False, CYAN, 'Cairo')],
    [(' Arabic, French, and English supported natively. Arabic detected at >30% Unicode Arabic characters; each language routed to its appropriate retriever.', 10, False, WHITE, 'Cairo')],
    [],
    [('Dual-Context Prompting', 10, False, CYAN, 'Cairo')],
    [(' Two distinct prompt blocks: [COURSE MATERIAL - RAG] for factual grounding (2,000 chars) and [STYLE ANCHORS] with 3–6 topic-relevant few-shot examples filtered by subject and Bloom level.', 10, False, WHITE, 'Cairo')],
    [],
    [('4-Axis Independent Evaluation', 10, False, CYAN, 'Cairo')],
    [(' Each question evaluated on: (1) Answer Extractability, (2) Bloom Level Alignment, (3) Question Quality, (4) Difficulty Fairness. Evaluator re-retrieves via question text to eliminate circular bias.', 10, False, WHITE, 'Cairo')],
])

# ── Slide 24 (Results – currently empty) – will be rebuilt as generation quality later

# ── Slide 26 (Models) – 3-axis → 4-axis ──────────────────────────────────────
slide26 = S[25]
tshapes26 = get_text_shapes(slide26)
models_box = [s for s in tshapes26 if 'deepseek' in s.text_frame.text.lower()][0]
replace_txBody(models_box, [
    [('deepseek/deepseek-v4-flash', 10, True, CYAN, 'Cairo Bold')],
    [(' Question generation (MCQ, SAQ, Exercise), temperature 0.75, structured JSON (Pydantic), multilingual, via OpenRouter.', 10, False, PWHITE, 'Cairo')],
    [],
    [('openai/gpt-oss-120b:free', 10, True, CYAN, 'Cairo Bold')],
    [(' 4-axis evaluation (Extractability, Bloom Level, Quality, Difficulty Fairness), temperature 0, strong reasoning, via OpenRouter.', 10, False, PWHITE, 'Cairo')],
    [],
    [('AraBERT (aubmindlab/bert-large-arabertv02)', 10, True, CYAN, 'Cairo Bold')],
    [(' Arabic embeddings (1,024-dim), FAISS indexing, handles morphology, via HuggingFace Inference API.', 10, False, PWHITE, 'Cairo')],
    [],
    [('BM25Okapi', 10, True, CYAN, 'Cairo Bold')],
    [(' Sparse retrieval + topic-aware few-shot sampling, fast (<5 ms), no inference cost.', 10, False, PWHITE, 'Cairo')],
    [],
    [('TF-IDF + LSA', 10, True, CYAN, 'Cairo Bold')],
    [(' Dense retrieval (FR/EN), 128-dim via SVD, no GPU required.', 10, False, PWHITE, 'Cairo')],
    [],
    [('FAISS (IndexFlatIP)', 10, True, CYAN, 'Cairo Bold')],
    [(' Cosine similarity search for all embeddings, CPU-based.', 10, False, PWHITE, 'Cairo')],
])

# ── Slide 27 (Thank You) – update text ───────────────────────────────────────
slide27 = S[26]
tshapes27 = get_text_shapes(slide27)
ty_box = [s for s in tshapes27 if 'Thank' in s.text_frame.text][0]
replace_txBody(ty_box, [
    [('Thank You', 30, False, WHITE, 'Zen Dots')],
    [],
    [('We sincerely thank our supervisor, ', 11, False, WHITE, 'Cairo'),
     ('Laid Kahloul', 11, True, CYAN, 'Cairo'),
     (', and our co-supervisor, ', 11, False, WHITE, 'Cairo'),
     ('Nour El Houda Ben Chaabene', 11, True, CYAN, 'Cairo'),
     (', for their invaluable guidance throughout this project.', 11, False, WHITE, 'Cairo')],
])

print("Step 1 done: existing slides modified.")

# ═══════════════════════════════════════════════════════════════════════════════
# STEP 2 – Create new slides (cloned, then text replaced)
# All new slides are appended at the end; we'll reorder in Step 3.
# Original slide count = 27 (indices 0-26).
# ═══════════════════════════════════════════════════════════════════════════════

# ── NEW: Research Gap & System Scope (clone from slide 6 style) ───────────────
ns_gap = clone_slide(prs, 5)  # clone Related Work slide style
tshapes = get_text_shapes(ns_gap)
title_box = [s for s in tshapes if 'Related Work' in s.text_frame.text or
             len(s.text_frame.text) < 50][0]
replace_txBody(title_box, [[('Research Gap & System Scope', 30, False, WHITE, 'Zen Dots')]])
content_box = max(tshapes, key=lambda s: len(s.text_frame.text))
replace_txBody(content_box, [
    [('Identified Gap', 10, True, CYAN, 'Cairo'),
     ('  No existing system simultaneously addresses bilingual IDP, language-aware hybrid retrieval, and independent post-generation quality evaluation — all required for real Algerian university archives.', 10, False, WHITE, 'Cairo')],
    [],
    [('Bilingual Archive Problem', 10, True, CYAN, 'Cairo'),
     ('  Algerian university archives consist of mixed French and Arabic PDFs with no machine-readable indexing. Standard NLP pipelines fail on bilingual, mixed-script corpora without dedicated routing and language-aware embeddings.', 10, False, WHITE, 'Cairo')],
    [],
    [('Existing Systems\' Limitations', 10, True, CYAN, 'Cairo'),
     ('  Related systems (Mucciaccia 2025, Isley 2025, Morris 2024) handle RAG grounding or schema-enforced generation, but none address bilingual corpora nor an independent Bloom-level classifier separate from the generator.', 10, False, WHITE, 'Cairo')],
    [],
    [('Our Proposed Scope', 10, True, CYAN, 'Cairo'),
     ('  Language-aware hybrid retrieval (BM25 + AraBERT/TF-IDF + RRF), structured generation with schema enforcement, and a 4-axis independent post-generation evaluator — without requiring model fine-tuning.', 10, False, WHITE, 'Cairo')],
])
IDX_RESEARCH_GAP = len(prs.slides) - 1  # 27

# ── NEW: Section divider — Related Work ───────────────────────────────────────
ns_rw_div = clone_slide(prs, 2)  # clone Problematic divider style
tshapes = get_text_shapes(ns_rw_div)
div_box = [s for s in tshapes if 'Problematic' in s.text_frame.text][0]
replace_txBody(div_box, [[('Related Work & Inspirations', 30, False, WHITE, 'Zen Dots')]])
IDX_RW_DIV = len(prs.slides) - 1  # 28

# ── NEW: Recent Key Studies ───────────────────────────────────────────────────
ns_recent = clone_slide(prs, 5)
tshapes = get_text_shapes(ns_recent)
title_box = [s for s in tshapes if len(s.text_frame.text.strip()) < 60 and s.text_frame.text.strip()][0]
replace_txBody(title_box, [[('Related Work — Recent Key Studies', 30, False, WHITE, 'Zen Dots')]])
content_box = max(tshapes, key=lambda s: len(s.text_frame.text))
replace_txBody(content_box, [
    [('Mucciaccia et al. (2025)', 10, True, CYAN, 'Cairo'),
     (' — Schema-enforced generation + automated LLM-as-judge review. Pydantic JSON output validation eliminates malformed responses. Schema enforcement raises factual accuracy to ~87% on domain-specific question banks.', 10, False, WHITE, 'Cairo')],
    [],
    [('Isley et al. (2025)', 10, True, CYAN, 'Cairo'),
     (' — RAG grounding + post-generation iterative refinement evaluated on 91 courses (~1,700 students). Retrieval-augmented generation substantially reduces hallucination vs. prompting alone.', 10, False, WHITE, 'Cairo')],
    [],
    [('Morris et al. (2024)', 10, True, CYAN, 'Cairo'),
     (' — Answer extractability as primary quality metric; introduced SAQ grading rubrics. A generated question is only valid if its answer can be unambiguously extracted from the source material.', 10, False, WHITE, 'Cairo')],
    [],
    [('Meissner et al. (2024)', 10, True, CYAN, 'Cairo'),
     (' — Independent Bloom classifier: do not trust the generator\'s self-reported cognitive level. Generators often mislabel levels — especially between Apply and Analyze.', 10, False, WHITE, 'Cairo')],
])
IDX_RECENT = len(prs.slides) - 1  # 29

# ── NEW: System Positioning table ─────────────────────────────────────────────
ns_pos = clone_slide(prs, 5)
tshapes = get_text_shapes(ns_pos)
title_box = [s for s in tshapes if len(s.text_frame.text.strip()) < 60 and s.text_frame.text.strip()][0]
replace_txBody(title_box, [[('System Positioning — Feature Comparison', 30, False, WHITE, 'Zen Dots')]])
content_box = max(tshapes, key=lambda s: len(s.text_frame.text))
replace_txBody(content_box, [[('No existing system addresses all of these simultaneously.', 11, False, CYAN, 'Cairo')]])

add_styled_table(
    ns_pos,
    left=350000, top=2200000, width=17500000, height=7000000,
    headers=['Feature', 'Mucciaccia\'25', 'Isley\'25', 'Morris\'24', 'Meissner\'24', 'This System'],
    rows=[
        ['Bilingual corpus (FR + AR)', '✗', '✗', '✗', '✗', '✓'],
        ['Language-aware hybrid retrieval', '✗', '✓', '✓', '✗', '✓'],
        ['Schema-enforced generation', '✓', '✓', '✗', '✗', '✓'],
        ['Post-generation multi-axis evaluator', '✓', '✓', '✗', '✗', '✓'],
        ['Independent Bloom-level classifier', '✗', '✗', '✗', '✓', '✓'],
        ['Topic-aware few-shot prompting', '✗', '✗', '✗', '✗', '✓'],
        ['Anti-duplication within session', '✗', '✗', '✗', '✗', '✓'],
    ],
    col_widths=[6500000, 2200000, 2200000, 2200000, 2200000, 2200000],
)
IDX_POSITIONING = len(prs.slides) - 1  # 30

# ── NEW: System Performance & Validation Results ──────────────────────────────
ns_perf = clone_slide(prs, 17)  # clone Contribution slide style
tshapes = get_text_shapes(ns_perf)
title_box = [s for s in tshapes if 'Contribution' in s.text_frame.text][0]
replace_txBody(title_box, [[('System Performance & Validation Results', 25, False, WHITE, 'Zen Dots')]])
content_box = [s for s in tshapes if 'Hybrid' in s.text_frame.text][0]
replace_txBody(content_box, [
    [('Generation Quality — 120 sampled questions across 5 subjects', 11, True, CYAN, 'Cairo')],
    [],
    [('87.5% factual accuracy (matches upper range of Mucciaccia et al.). Law scores lowest due to sparse Arabic corpus (66 questions, 11 JSON files).', 10, False, WHITE, 'Cairo')],
    [],
    [('Retrieval Effectiveness — 50 held-out queries', 11, True, CYAN, 'Cairo')],
    [],
    [('Best config: BM25 + AraBERT/MiniLM (RRF + routing) → 85.0% French, 81.0% Arabic, 83.8% overall.', 10, False, WHITE, 'Cairo')],
    [('Language-aware routing to AraBERT: +10 points on Arabic queries (81.0% vs 71.0%).', 10, False, WHITE, 'Cairo')],
    [],
    [('Dataset Scale', 11, True, CYAN, 'Cairo')],
    [('137 real bilingual university documents → 86 structured JSON files → 241 exercises → 704 Bloom-labelled questions — without model fine-tuning.', 10, False, WHITE, 'Cairo')],
])
IDX_PERF = len(prs.slides) - 1  # 31

# ── NEW: Section divider — System Architecture ────────────────────────────────
ns_arch_div = clone_slide(prs, 2)
tshapes = get_text_shapes(ns_arch_div)
div_box = [s for s in tshapes if 'Problematic' in s.text_frame.text][0]
replace_txBody(div_box, [[('System Architecture', 30, False, WHITE, 'Zen Dots')]])
IDX_ARCH_DIV = len(prs.slides) - 1  # 32

# ── NEW: Four-Layer Architecture ──────────────────────────────────────────────
ns_4layer = clone_slide(prs, 17)
tshapes = get_text_shapes(ns_4layer)
title_box = [s for s in tshapes if 'Contribution' in s.text_frame.text][0]
replace_txBody(title_box, [[('Four-Layer Architecture', 30, False, WHITE, 'Zen Dots')]])
content_box = [s for s in tshapes if 'Hybrid' in s.text_frame.text][0]
replace_txBody(content_box, [
    [('Layer 1 — Raw Document Corpus', 11, True, CYAN, 'Cairo')],
    [(' 137 files: PDF · DOCX · JPEG  |  5 subjects  |  French & Arabic.  PDF upload via web → docparser_service launches background parse.', 10, False, WHITE, 'Cairo')],
    [],
    [('Layer 2 — Intelligent Document Processing (docParser)', 11, True, CYAN, 'Cairo')],
    [(' S1: PDF→Markdown  |  S2: LLM + Instructor  |  S3: Grounding Validation  |  S4: Bloom Classification  |  S5: VLM Diagram Transcription.', 10, False, WHITE, 'Cairo')],
    [(' Output → /api/index-file (cloud RAG) + /api/invalidate-cache', 10, False, CYAN, 'Cairo')],
    [],
    [('Layer 3 — RAG-Augmented Generation Engine (exam-forge)', 11, True, CYAN, 'Cairo')],
    [(' BM25 + FAISS + RRF Retrieval  |  DeepSeek Generator  |  4-axis Evaluator  |  TemplateAnalyzer', 10, False, WHITE, 'Cairo')],
    [(' Exposed via REST API / SSE streaming.', 10, False, CYAN, 'Cairo')],
    [],
    [('Layer 4 — Web Application', 11, True, CYAN, 'Cairo')],
    [(' FastAPI + SQLite  |  JWT / Google OAuth 2.0  |  React + Vite  |  SSE real-time streaming.', 10, False, WHITE, 'Cairo')],
])
IDX_4LAYER = len(prs.slides) - 1  # 33

# ── NEW: docParser 5-Stage Pipeline ──────────────────────────────────────────
ns_parser = clone_slide(prs, 9)  # clone Folder Scanning style
tshapes = get_text_shapes(ns_parser)
# Find the largest text box for content
content_boxes = sorted([s for s in tshapes if s.has_text_frame], key=lambda s: -len(s.text_frame.text))
if len(content_boxes) >= 2:
    replace_txBody(content_boxes[1], [
        [('Stage 1 — PDF → Markdown', 10, False, CYAN, 'Canva Sans')],
        [(' Docling with TableFormer (table structure) + RTL-safe export. PyMuPDF hybrid fill for formula/pseudocode recovery.', 10, False, WHITE, 'Canva Sans')],
        [('Stage 2 — LLM Extraction with Self-Healing Schema', 10, False, CYAN, 'Canva Sans')],
        [(' instructor intercepts Pydantic errors → correction prompts (max_retries=3). Schema: Exam ⊃ Exercise ⊃ Question.', 10, False, WHITE, 'Canva Sans')],
        [('Stage 3 — Grounding Validation', 10, False, CYAN, 'Canva Sans')],
        [(' Confidence score c = max(0, 1 − 0.15|W| − 0.5·1[E=∅] − …). Flag if c < 0.70.', 10, False, WHITE, 'Canva Sans')],
        [('Stage 4 — Bloom Classification', 10, False, CYAN, 'Canva Sans')],
        [(' Independent classifier (gpt-oss-120b, temp=0): Factual / Conceptual / Procedural / Metacognitive.', 10, False, WHITE, 'Canva Sans')],
        [('Stage 5 — VLM Diagram Transcription', 10, False, CYAN, 'Canva Sans')],
        [(' Qwen2.5-VL-72B (primary) or Gemini-2.0-Flash (fallback); linked to nearest preceding question.', 10, False, WHITE, 'Canva Sans')],
    ])
    replace_txBody(content_boxes[0], [[('docParser: 5-Stage IDP Pipeline', 30, False, WHITE, 'Zen Dots')]])
IDX_PARSER = len(prs.slides) - 1  # 34

# ── NEW: docParser Processing Results ────────────────────────────────────────
ns_parser_res = clone_slide(prs, 5)
tshapes = get_text_shapes(ns_parser_res)
title_box2 = [s for s in tshapes if len(s.text_frame.text.strip()) < 60 and s.text_frame.text.strip()][0]
replace_txBody(title_box2, [[('docParser Processing Results (137 documents)', 25, False, WHITE, 'Zen Dots')]])
content_box2 = max(tshapes, key=lambda s: len(s.text_frame.text))
replace_txBody(content_box2, [[('Of 33 rejected: 21 answer-key sheets (no question text) + 12 low-resolution scans.', 10, False, CYAN, 'Cairo')]])

add_styled_table(
    ns_parser_res,
    left=350000, top=2000000, width=17500000, height=7200000,
    headers=['Metric', 'Count', 'Proportion'],
    rows=[
        ['Successful Docling conversion', '124', '90.5%'],
        ['PyMuPDF fallback used', '13', '9.5%'],
        ['Hybrid formula/code fill triggered', '41', '33.1%'],
        ['Accepted (c ≥ 0.70)', '86', '62.8%'],
        ['Flagged for review (c < 0.70)', '18', '13.1%'],
        ['Rejected (no extractable questions)', '33', '24.1%'],
        ['Visual diagrams extracted (Stage 5)', '47', '—'],
        ['Described by Qwen2.5-VL-72B', '42', '89.4%'],
        ['Described by Gemini-2.0-Flash', '5', '10.6%'],
    ],
    col_widths=[11000000, 3200000, 3300000],
)
IDX_PARSER_RES = len(prs.slides) - 1  # 35

# ── NEW: Web Application Overview ────────────────────────────────────────────
ns_webapp = clone_slide(prs, 17)
tshapes = get_text_shapes(ns_webapp)
title_box = [s for s in tshapes if 'Contribution' in s.text_frame.text][0]
replace_txBody(title_box, [[('Web Application Overview', 30, False, WHITE, 'Zen Dots')]])
content_box = [s for s in tshapes if 'Hybrid' in s.text_frame.text][0]
replace_txBody(content_box, [
    [('Authentication', 10, True, CYAN, 'Cairo'),
     ('  Login via email/password or Google OAuth 2.0. Light mode and dark mode supported.', 10, False, WHITE, 'Cairo')],
    [],
    [('Subject Dashboard', 10, True, CYAN, 'Cairo'),
     ('  Main page shows all subject modules with course cards. Each card provides access to full lesson details, exploration, and module management.', 10, False, WHITE, 'Cairo')],
    [],
    [('Exam Builder', 10, True, CYAN, 'Cairo'),
     ('  Select subject, configure question type (MCQ / SAQ / Exercise) and Bloom level, then launch generation. Questions stream card-by-card via SSE in real time.', 10, False, WHITE, 'Cairo')],
    [],
    [('Archive (My Archive)', 10, True, CYAN, 'Cairo'),
     ('  Full browsable history of all generation sessions. Includes Favorites, Trash (negative examples for system teaching), and HISSTORIQ (complete generation log).', 10, False, WHITE, 'Cairo')],
    [],
    [('Exam Builder — Assemble & Export', 10, True, CYAN, 'Cairo'),
     ('  Step 1: Save exercises to favorites. Step 2: Drag and build a draft document. Step 3: Export to PDF. Step 4: Print.', 10, False, WHITE, 'Cairo')],
    [],
    [('Tech Stack', 10, True, CYAN, 'Cairo'),
     ('  FastAPI + SQLite  |  React + Vite  |  JWT + Google OAuth 2.0  |  SSE real-time streaming.', 10, False, WHITE, 'Cairo')],
])
IDX_WEBAPP = len(prs.slides) - 1  # 36

# ── NEW: Exam Generation Workflow ─────────────────────────────────────────────
ns_workflow = clone_slide(prs, 17)
tshapes = get_text_shapes(ns_workflow)
title_box = [s for s in tshapes if 'Contribution' in s.text_frame.text][0]
replace_txBody(title_box, [[('Exam Generation Workflow', 30, False, WHITE, 'Zen Dots')]])
content_box = [s for s in tshapes if 'Hybrid' in s.text_frame.text][0]
replace_txBody(content_box, [
    [('1. Access Subject', 11, True, CYAN, 'Cairo'),
     ('  — User selects a subject card to enter its dashboard.', 10, False, WHITE, 'Cairo')],
    [('2. View Subject Details', 11, True, CYAN, 'Cairo'),
     ('  — Dashboard shows course materials and available options.', 10, False, WHITE, 'Cairo')],
    [('3. Add Course Content', 11, True, CYAN, 'Cairo'),
     ('  — User enters title and description for new material.', 10, False, WHITE, 'Cairo')],
    [('4. Access AI Generator', 11, True, CYAN, 'Cairo'),
     ('  — User selects the generator button in the top right.', 10, False, WHITE, 'Cairo')],
    [('5. AI Settings', 11, True, CYAN, 'Cairo'),
     ('  — Choose question type (MCQ / SAQ / Exercise) and Bloom level, then launch generation.', 10, False, WHITE, 'Cairo')],
    [('6. Generating', 11, True, CYAN, 'Cairo'),
     ('  — Questions stream card-by-card via SSE in real time.', 10, False, WHITE, 'Cairo')],
    [('7. Review', 11, True, CYAN, 'Cairo'),
     ('  — Section cards populate with results: Mixed level · MCQ · SAQ sections.', 10, False, WHITE, 'Cairo')],
    [('8. Accept / Edit / Reject', 11, True, CYAN, 'Cairo'),
     ('  — The professor validates, edits, or rejects each question before finalizing.', 10, False, WHITE, 'Cairo')],
])
IDX_WORKFLOW = len(prs.slides) - 1  # 37

# ── NEW: Section divider — Models & Technical Stack ───────────────────────────
ns_models_div = clone_slide(prs, 22)  # clone Results divider style
tshapes = get_text_shapes(ns_models_div)
div_box = [s for s in tshapes if 'Results' in s.text_frame.text][0]
replace_txBody(div_box, [[('Models & Technical Stack', 30, False, WHITE, 'Zen Dots')]])
IDX_MODELS_DIV = len(prs.slides) - 1  # 38

# ── Rebuild slide 24 (Results) as Generation Quality table ───────────────────
slide24 = S[23]
tshapes24 = get_text_shapes(slide24)
res_title = [s for s in tshapes24 if 'Results' in s.text_frame.text]
if res_title:
    replace_txBody(res_title[0], [[('Generation Quality (120 sampled questions)', 30, False, WHITE, 'Zen Dots')]])

add_styled_table(
    slide24,
    left=350000, top=2000000, width=17500000, height=5500000,
    headers=['Subject', 'Factual Accuracy', 'Bloom Alignment', 'Quality: Good', 'Difficulty Fair'],
    rows=[
        ['Algorithms',     '87.5%', '79.2%', '75.0%', '83.3%'],
        ['Software Eng.',  '91.7%', '83.3%', '83.3%', '87.5%'],
        ['Compilation',    '87.5%', '83.3%', '79.2%', '83.3%'],
        ['Commerce',       '87.5%', '79.2%', '79.2%', '87.5%'],
        ['Law',            '83.3%', '75.0%', '70.8%', '79.2%'],
        ['Overall',        '87.5%', '79.2%', '77.5%', '84.2%'],
    ],
)
# Add note below table
add_textbox(slide24,
    left=350000, top=7700000, width=17500000, height=600000,
    paragraphs=[[('87.5% factual accuracy matches the upper range of Mucciaccia et al. Law scores lowest: small Arabic corpus (66 questions, 11 JSON files).', 10, False, CYAN, 'Cairo')]],
    word_wrap=True)

# ── NEW: Retrieval Effectiveness table ───────────────────────────────────────
ns_retrieval = clone_slide(prs, 22)  # clone Results section-divider for layout
tshapes = get_text_shapes(ns_retrieval)
div_box = [s for s in tshapes if 'Results' in s.text_frame.text][0]
replace_txBody(div_box, [[('Retrieval Effectiveness (50 held-out queries)', 25, False, WHITE, 'Zen Dots')]])

add_styled_table(
    ns_retrieval,
    left=350000, top=2000000, width=17500000, height=5500000,
    headers=['Configuration', 'French Subjects', 'Arabic (Law)', 'Overall'],
    rows=[
        ['BM25 only',                           '79.0%', '62.0%', '74.0%'],
        ['Dense only — MiniLM',                 '81.0%', '68.0%', '77.2%'],
        ['Dense only — AraBERT',                '83.0%', '81.0%', '82.6%'],
        ['BM25 + MiniLM (RRF, no routing)',     '84.0%', '71.0%', '80.8%'],
        ['BM25 + AraBERT/MiniLM (RRF+routing)', '85.0%', '81.0%', '83.8%'],
    ],
    col_widths=[8000000, 3200000, 3200000, 3100000],
)
add_textbox(ns_retrieval,
    left=350000, top=7700000, width=17500000, height=600000,
    paragraphs=[[('Language-aware routing to AraBERT: +10 points on Arabic queries (81.0% vs 71.0%) — confirms necessity of dedicated Arabic embeddings.', 10, False, CYAN, 'Cairo')]],
    word_wrap=True)
IDX_RETRIEVAL = len(prs.slides) - 1  # 39

print("Step 2 done: new slides created.")

# ═══════════════════════════════════════════════════════════════════════════════
# STEP 3 – Reorder slides to final sequence
# ═══════════════════════════════════════════════════════════════════════════════
# Current indices after all additions:
# 0-26:  original slides (some modified)
# 27: Research Gap (IDX_RESEARCH_GAP)
# 28: RW section divider (IDX_RW_DIV)
# 29: Recent Key Studies (IDX_RECENT)
# 30: System Positioning (IDX_POSITIONING)
# 31: System Perf (IDX_PERF)
# 32: Arch section divider (IDX_ARCH_DIV)
# 33: Four-Layer (IDX_4LAYER)
# 34: docParser pipeline (IDX_PARSER)
# 35: docParser results (IDX_PARSER_RES)
# 36: Web App (IDX_WEBAPP)
# 37: Workflow (IDX_WORKFLOW)
# 38: Models section divider (IDX_MODELS_DIV)
# 39: Retrieval table (IDX_RETRIEVAL)
#
# Old slides to REMOVE after reorder: 7 (Pipeline), 8 (Offline), 11 (Online), 23 (old empty Results)
# Wait - slide 23 (Results section divider) we KEEP; slide 23 (0-based index 22) is the Results divider.
# Original slide 24 (idx 23) was empty Results - we've now filled it with the table.
# But we're keeping it in the flow.
#
# Desired final order (using CURRENT indices before any moves):
# 0  → slide  1 Title
# 1  → slide  2 Contents
# 2  → slide  3 Section Problematic
# 3  → slide  4 Problematic
# 4  → slide  5 Objectives
# 27 → slide  6 Research Gap (NEW)
# 28 → slide  7 Section Related Work (NEW)
# 5  → slide  8 Related Work General
# 29 → slide  9 Recent Key Studies (NEW)
# 30 → slide 10 System Positioning (NEW)
# 6  → slide 11 Shared Pillars
# 16 → slide 12 Section Proposition (orig slide 17)
# 17 → slide 13 Contribution (orig slide 18)
# 31 → slide 14 System Performance (NEW)
# 32 → slide 15 Section Architecture (NEW)
# 33 → slide 16 Four-Layer (NEW)
# 9  → slide 17 Doc Ingestion (orig slide 10)
# 34 → slide 18 docParser Pipeline (NEW)
# 35 → slide 19 docParser Results (NEW)
# 10 → slide 20 Indexing (orig slide 11)
# 12 → slide 21 Topic-Aware (orig slide 13)
# 13 → slide 22 RAG Context (orig slide 14)
# 14 → slide 23 Structured Gen (orig slide 15)
# 15 → slide 24 Post-Gen Eval (orig slide 16)
# 36 → slide 25 Web App (NEW)
# 37 → slide 26 Workflow (NEW)
# 18 → slide 27 Section Problems (orig slide 19)
# 19 → slide 28 Problems 1&2 (orig slide 20)
# 20 → slide 29 Problems 3&4 (orig slide 21)
# 21 → slide 30 Problems 5&6 (orig slide 22)
# 22 → slide 31 Section Results (orig slide 23)
# 23 → slide 32 Gen Quality table (orig slide 24, modified)
# 39 → slide 33 Retrieval table (NEW)
# 24 → slide 34 Technical Achievements (orig slide 25)
# 38 → slide 35 Section Models (NEW)
# 25 → slide 36 Models & Tech (orig slide 26)
# 26 → slide 37 Thank You (orig slide 27)
# SKIP: 7 (Pipeline), 8 (Offline), 11 (Online) → don't include

desired_order = [
    0, 1, 2, 3, 4,           # Title, Contents, Sect Prob, Problematic, Objectives
    27, 28,                   # Research Gap, Sect Related Work
    5, 29, 30, 6,             # Related Work, Recent Studies, Positioning, Shared Pillars
    16, 17, 31,               # Sect Proposition, Contribution, System Perf
    32, 33,                   # Sect Architecture, Four-Layer
    9, 34, 35, 10,            # Doc Ingestion, docParser Pipeline, docParser Results, Indexing
    12, 13, 14, 15,           # Topic-Aware, RAG, Gen, Post-Gen Eval
    36, 37,                   # Web App, Workflow
    18, 19, 20, 21,           # Sect Problems, Prob 1&2, 3&4, 5&6
    22, 23, 39, 24,           # Sect Results, Gen Quality, Retrieval, Technical Achievements
    38, 25, 26,               # Sect Models, Models & Tech, Thank You
]

# Build new sldIdLst order
xml_slides = prs.slides._sldIdLst
all_ids = list(xml_slides)

# Remove all from list
for el in all_ids:
    xml_slides.remove(el)

# Re-insert in desired order
for idx in desired_order:
    xml_slides.append(all_ids[idx])

print("Step 3 done: slides reordered.")
print(f"Final slide count: {len(prs.slides)}")

prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
