#!/usr/bin/env python3
"""Fix specific issues in the updated PPTX."""
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

INPUT  = 'AI-powered assessment platform for higher education (updated).pptx'
OUTPUT = 'AI-powered assessment platform for higher education (updated).pptx'

WHITE  = RGBColor(0xF2, 0xF2, 0xF2)
PWHITE = RGBColor(0xFF, 0xFF, 0xFF)
CYAN   = RGBColor(0x77, 0xE8, 0xFF)
PINK   = RGBColor(0xF5, 0x9A, 0xF4)

def _make_run(text, size_pt, bold, color: RGBColor, font_name):
    r = etree.Element(qn('a:r'))
    rPr = etree.SubElement(r, qn('a:rPr'))
    rPr.set('lang', 'en-US')
    rPr.set('dirty', '0')
    if size_pt:
        rPr.set('sz', str(int(size_pt * 100)))
    if bold:
        rPr.set('b', '1')
    if color:
        sf = etree.SubElement(rPr, qn('a:solidFill'))
        sc = etree.SubElement(sf, qn('a:srgbClr'))
        sc.set('val', str(color).upper())
    if font_name:
        lat = etree.SubElement(rPr, qn('a:latin'))
        lat.set('typeface', font_name)
    t = etree.SubElement(r, qn('a:t'))
    t.text = text
    return r

def replace_txBody(shape, paragraphs):
    txBody = shape.text_frame._txBody
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)
    for para in paragraphs:
        p_elem = etree.Element(qn('a:p'))
        if not para:
            txBody.append(p_elem)
            continue
        for text, sz, bold, color, font in para:
            p_elem.append(_make_run(text, sz, bold, color, font))
        txBody.append(p_elem)

def clear_shape_text(shape):
    """Remove all text from a shape's text frame."""
    txBody = shape.text_frame._txBody
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)
    txBody.append(etree.Element(qn('a:p')))

def restore_slide_number(shape):
    """Restore a slide-number placeholder to its default auto field."""
    txBody = shape.text_frame._txBody
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)
    # Add empty paragraph (slide number will auto-populate in PPT)
    p = etree.Element(qn('a:p'))
    fld = etree.SubElement(p, qn('a:fld'))
    fld.set('{http://schemas.openxmlformats.org/drawingml/2006/main}id',
            '{00000000-0000-0000-0000-000000000000}')
    fld.set('type', 'slidenum')
    rPr = etree.SubElement(fld, qn('a:rPr'))
    rPr.set('lang', 'en-US')
    t = etree.SubElement(fld, qn('a:t'))
    t.text = '‹#›'
    txBody.append(p)

prs = Presentation(INPUT)
S = prs.slides

# ── Fix Slide 6 (Research Gap) ─────────────────────────────────────────────
# Freeform 12 accidentally got the title — clear it
# TextBox 15 still says "Related Work" — set it to proper title
slide6 = S[5]
for s in slide6.shapes:
    if s.name == 'Freeform 12' and s.has_text_frame:
        clear_shape_text(s)
    if s.name == 'TextBox 15' and s.has_text_frame:
        replace_txBody(s, [[('Research Gap & System Scope', 30, False, WHITE, 'Zen Dots')]])

# ── Fix Slide 7 (Section: Related Work divider) ────────────────────────────
# Page number placeholder shows "3" (from source) — restore to auto
slide7 = S[6]
for s in slide7.shapes:
    if 'réservé' in s.name.lower() and s.has_text_frame:
        restore_slide_number(s)

# ── Fix Slide 18 (docParser Pipeline) ─────────────────────────────────────
# Slide-number placeholder got stage content — restore it
# TextBox 22 is empty — put stages content there
slide18 = S[17]
for s in slide18.shapes:
    if 'réservé' in s.name.lower() and s.has_text_frame:
        restore_slide_number(s)
    if s.name == 'TextBox 22' and s.has_text_frame:
        replace_txBody(s, [
            [('Stage 1 — PDF → Markdown', 10, False, CYAN, 'Canva Sans')],
            [(' Docling with TableFormer (table structure) + RTL-safe export. PyMuPDF hybrid fill for formula/pseudocode recovery.', 10, False, WHITE, 'Canva Sans')],
            [('Stage 2 — LLM Extraction with Self-Healing Schema', 10, False, CYAN, 'Canva Sans')],
            [(' instructor intercepts Pydantic errors → correction prompts (max_retries=3). Schema: Exam ⊃ Exercise ⊃ Question.', 10, False, WHITE, 'Canva Sans')],
            [('Stage 3 — Grounding Validation', 10, False, CYAN, 'Canva Sans')],
            [(' Confidence score c = max(0, 1 − 0.15|W| − 0.5·1[E=∅] − …). Flag if c < 0.70.', 10, False, WHITE, 'Canva Sans')],
            [('Stage 4 — Bloom Classification', 10, False, CYAN, 'Canva Sans')],
            [(' Independent classifier (gpt-oss-120b, temp=0): Factual / Conceptual / Procedural / Metacognitive.', 10, False, WHITE, 'Canva Sans')],
            [('Stage 5 — VLM Diagram Transcription', 10, False, CYAN, 'Canva Sans')],
            [(' Qwen2.5-VL-72B (primary) or Gemini-2.0-Flash (fallback); linked to nearest preceding question.', 10, False, WHITE, 'Canva Sans')],
        ])
    # Also clear the old "Folder Scanning" content from TextBox 23 which now holds the title
    # (it already says docParser: 5-Stage IDP Pipeline from the update script — keep it)

# ── Fix placeholder slide numbers on all cloned slides ─────────────────────
# Slides that were cloned may have wrong hard-coded numbers
# Cloned section dividers: 7(6), 15(14), 27(26), 31(30), 33(32), 35(34)
# Cloned content slides: 6(5), 9(8), 10(9), 14(13), 16(15), 18(17), 19(18),
#                         25(24), 26(25), 33(32)
for idx in [5, 6, 8, 9, 13, 14, 17, 18, 24, 25, 29, 30, 32, 33, 34, 37, 38]:
    if idx >= len(S):
        continue
    slide = S[idx]
    for s in slide.shapes:
        if 'réservé' in s.name.lower() and s.has_text_frame:
            restore_slide_number(s)

prs.save(OUTPUT)
print(f"Fixed and saved: {OUTPUT}")
